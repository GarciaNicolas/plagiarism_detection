import os
import shutil
import time
import ast

from helper import pd, re, np, author_synonyms
from unicodedata import normalize
from doc2docx import convert as doc2docx
from pptx import Presentation
from docx import Document
from docx.opc.constants import RELATIONSHIP_TYPE as RT
from pdf2docx import Converter
from bs4 import BeautifulSoup
from urllib.request import urlopen, Request
from processor import get_corpus,process_indexed_corpus


#---Read-----------------------------------------------------------------------------------------------------------------------
def read_file(path):
    """
    Read a file and return a dataframe with the file name, the text and the citations
    """
    df = pd.DataFrame(columns=['filename','text','citations','headers'])

    #Read the file depending on its extension and concatenate its content to the dataframe
    if path.endswith('.docx'):
        df = read_docx(path)
    elif path.endswith('.pptx'):
        df = read_pptx(path)
    elif path.endswith('.doc') or path.endswith('.pdf'):
        df = read_and_convert_file(path, df)
    else:
        raise Exception('FILE EXTENSION NOT SUPPORTED: Only support .docx, .doc, .pdf and .pptx files')

    return df

def read_path(path):
    """
    Read all files in a path and return a dataframe with the file name, the text and the citations
    """
    files = os.listdir(path)

    if len(files) == 1:
        return read_file(path + files[0])

    df = pd.DataFrame(columns=['filename','text','citations','headers'])

    #Create a list of .doc files and concatenate its content to the dataframe
    docs_and_docxs = [file for file in files if file.endswith('.doc') or file.endswith('.pdf')]
    df = read_pdf_docs(docs_and_docxs, df, path)

    #Create a list of files which are not .doc
    not_docs_or_docxs = [file for file in files if file not in docs_and_docxs]

    #Read each file depending on its extension and concatenate its content to the dataframe
    for file in not_docs_or_docxs:

        file_path = path + file
        if file.endswith('.docx'):
            file_df = read_docx(file_path)
        elif file.endswith('.pptx'):
            file_df = read_pptx(file_path)
        else:
            raise Exception('FILE EXTENSION NOT SUPPORTED: Only support .docx, .doc, .pdf and .pptx files')
        
        df = pd.concat([df,file_df], ignore_index=True)
    return df

def read_database():
    """
    Read the database, fix the lists and concatenate it to the new dataframe
    """
    db = pd.read_csv('../data.csv')
    db = get_lists(db)

    return db

def read_urls(urls):
    """
    Read the text of each hyperlink and concatenate it to the dataframe
    """
    df = pd.DataFrame(columns=['url','corpus', 'processed_corpus'])
    for url in urls:
            try:
                req = Request(
                            url=url,
                            headers={'User-Agent': 'Mozilla/5.0'})
                html = urlopen(req).read()
            except:
                continue
            soup = BeautifulSoup(html, features="html.parser")

            #Kill all script and style elements
            for script in soup(["script", "style"]):
                script.extract()

            #Get text
            text = soup.get_text()

            #Break into lines and remove leading and trailing space on each
            lines = (line.strip() for line in text.splitlines())
            
            #Break multi-headlines into a line each
            chunks = (phrase.strip() for line in lines for phrase in line.split("  "))

            #Drop blank lines
            text = '. '.join(chunk for chunk in chunks if chunk)
            corpus = get_corpus(text)
            processed_corpus = process_indexed_corpus(corpus)
            

            df = pd.concat([df, pd.DataFrame({'url':url,'corpus':[corpus], 'processed_corpus': [processed_corpus]}, index=[0])], ignore_index=True)
    return df    

#---Type of file---------------------------------------------------------------------------------------------------------------

def read_docx(file):
    """
    Read a docx file and return a dataframe with the file name, the text and the citations
    """
    doc = Document(file)
    hyperlinks = []
    text = ""

    headers = get_headers(doc)
    footers = get_footers(doc)

    names = names_in_table(doc.tables)
    hyperlinks += (docx_hyperlinks(doc))
    for paragraph in doc.paragraphs:
        paragraph_text = paragraph.text
        paragraph_text = remove_section(paragraph_text, headers)
        paragraph_text = remove_section(paragraph_text, footers)
        hyperlinks += text_hyperlinks(paragraph_text)
        text += paragraph_text + ' \n '
    text = clean_special_characters(text)

    headers = [clean_special_characters(header) for header in headers]

    hyperlinks = cleaning_citations(hyperlinks)
    df = pd.DataFrame({'filename': path_to_filename(file), 'text': [text],'author':[names] if names else np.nan, 'citations': [hyperlinks] if hyperlinks else np.nan, 'headers': [headers] if headers else np.nan})
    return df

def read_pptx(path):
    """
    Read a pptx file and return a dataframe with the file name, the text and the citations
    """
    text = ""
    hyperlinks = []

    #Reading the document by slides 
    for slide in Presentation(path).slides:
        #Looking for shapes which have text and extracting text and hyperlinks from them
        for shape in slide.shapes:
            if not shape.has_text_frame:
                continue
            for paragraph in shape.text_frame.paragraphs:
                    for run in paragraph.runs:
                        
                        #Looking for hyperlinks in the text
                        hyperlinks = hyperlinks + text_hyperlinks(run.text)
                        
                        #Cleaning special characters
                        text = text +'. ' + run.text
    hyperlinks = list(set(hyperlinks))
    df = pd.DataFrame({'filename': path_to_filename(path), 'text': [text], 'citations': [hyperlinks] if hyperlinks else None})

    return df

def read_pdf_docs(docs, df, path):
    """
    Receive a list of .doc and .pdf files, convert them to .docx files and read them,
    then concatenate the new dataframes of each new .docx file into a dataframe and return it
    """

    dot_docs = [doc for doc in docs if doc.endswith('.doc')]
    dot_pdfs = [doc for doc in docs if doc.endswith('.pdf')]

    #If there is a temporal folder, remove it
    if os.path.isdir('./data_temp'):

        #Remove the temporal folder and its files
        shutil.rmtree('./data_temp')

    os.mkdir('./data_temp/')    
    os.mkdir('./data_temp/docs_files')
    os.mkdir('./data_temp/pdfs_files')
    os.mkdir('./data_temp/docxs_files')

    #Move files into its temporal folders
    for doc in dot_docs:
        shutil.copy(path + doc, './data_temp/docs_files/')
    for doc in dot_pdfs:
        shutil.copy(path + doc, './data_temp/pdfs_files/')

    try:

        #Converting .doc files to .docx files and then to .pdf files
        doc2docx('./data_temp/docs_files/', './data_temp/docxs_files/')
        time.sleep(.5)
        pdf2docx(dot_pdfs)
        time.sleep(.5)

        #Read each .pdf file and concatenate the dataframes
        for file in os.listdir('./data_temp/docxs_files/'):
            file_df = read_docx('./data_temp/docxs_files/'+file)
            df = pd.concat([df,file_df], ignore_index=True)

    except:

        #Remove the temporal folder and its files
        shutil.rmtree('./data_temp')
        raise Exception('ERROR: Could not convert files to .pdf files and read them')

    #Remove the temporal folder and its files
    shutil.rmtree('./data_temp')

    return df

#---Convert--------------------------------------------------------------------------------------------------------------------

def read_and_convert_file(path, df):
    """
    Read and convert only one file with .doc or .pdf format
    """
    
    if path.endswith('.doc'):
        doc2docx(path, './temp.docx')
    elif path.endswith('.pdf'):
        cv = Converter(path)
        cv.convert('./temp.pdf')
        cv.close()
    else:
        raise Exception('ERROR: File format not supported')

    file_df = read_docx('./temp.docx')
    df = pd.concat([df,file_df], ignore_index=True)
    os.remove('./temp.docx')

    return df

def read_and_convert_from_path(docs, df, path):
    """
    Receive a list of .doc and .pdf files, convert them to .docx files and read them,
    then concatenate the new dataframes of each new .docx file into a dataframe and return it
    """

    dot_docs = [doc for doc in docs if doc.endswith('.doc')]
    dot_pdfs = [doc for doc in docs if doc.endswith('.pdf')]

    #If there is a temporal folder, remove it
    if os.path.isdir('./data_temp'):

        #Remove the temporal folder and its files
        shutil.rmtree('./data_temp')

    os.mkdir('./data_temp/')    
    os.mkdir('./data_temp/docs_files')
    os.mkdir('./data_temp/pdfs_files')
    os.mkdir('./data_temp/docxs_files')

    #Copy files into temporal folders
    for doc in dot_docs:
        shutil.copy(path + doc, './data_temp/docs_files/')
    for doc in dot_pdfs:
        shutil.copy(path + doc, './data_temp/pdfs_files/')

    try:

        #Converting .doc files to .docx files and then to .pdf files
        doc2docx('./data_temp/docs_files/', './data_temp/docxs_files/')
        time.sleep(.5)
        pdf2docx(dot_pdfs)
        time.sleep(.5)

        #Read each .pdf file and concatenate the dataframes
        for file in os.listdir('./data_temp/docxs_files/'):
            file_df = read_docx('./data_temp/docxs_files/'+file)
            df = pd.concat([df,file_df], ignore_index=True)

    except:

        #Remove the temporal folder and its files
        shutil.rmtree('./data_temp')
        raise Exception('ERROR: Could not convert files to .pdf files and read them')

    #Remove the temporal folder and its files
    shutil.rmtree('./data_temp')

    return df

def pdf2docx(pdf_files):
    filenames = [os.path.splitext(doc)[0] for doc in pdf_files]
    path_docx = './data_temp/docxs_files/'
    path_pdf = './data_temp/pdfs_files/'
    for file in filenames:
        cv = Converter(path_pdf + file + '.pdf')
        cv.convert(path_docx + file + '.docx', multi_processing=True)
        cv.close()

#---Clean----------------------------------------------------------------------------------------------------------------------

def cleaning_citations(strings):
    strings = list(set(strings))
    
    strings = [string for string in strings if not string.startswith('mailto:')]
    return strings

def clean_special_characters(string):
    """
    Receives a string and returns it without those specials characters which are troublesome
    """
    string = string.replace('\u200b',' ')
    string = string.replace('\xad',' ').replace('\xa0',' ')
    string = string.replace('\t',' ')
    string = re.sub('(\s*\n\s*)+', ' \n ', string)
    
    #Remove accents
    string = re.sub(
        r"([^n\u0300-\u036f]|n(?!\u0303(?![\u0300-\u036f])))[\u0300-\u036f]+", r"\1", 
        normalize( "NFD", string), 0, re.I)
    
    return string

def remove_section(string, list_sections):
    """
    Receives a string and a list of strings belonging to a section and returns the string without these section
    """
    for section in list_sections:
        string = string.replace(section,'')
    return string

def get_lists(df):
    """
    Converts lists represented as strings to lists
    """
    df.loc[df.citations.notnull(), 'citations'] = df[df.citations.notnull()]['citations'].apply(ast.literal_eval)
    df.loc[df.author.notnull(), 'author'] = df[df.author.notnull()]['author'].apply(ast.literal_eval)
    df.loc[df.headers.notnull(), 'headers'] = df[df.headers.notnull()]['headers'].apply(ast.literal_eval)
    df['processed_corpus'] = df['processed_corpus'].apply(ast.literal_eval)

    #Chek if it's not empty?? I really don't know what I tried to do here :/
    df.loc[df.headers.notnull(), 'headers']=df[df.headers.notnull()].headers.apply(lambda x: [header for header in x if header])

    return df

#---Extract--------------------------------------------------------------------------------------------------------------------

def text_hyperlinks(string):
    """
    Return a list of hyperlinks in a string
    """

    #Regular expression of hyperlinks
    regex = r"(https?://\S+)"

    #Find all hyperlinks in a string
    hyperlinks = re.findall(regex,string)
    
    return hyperlinks

def docx_hyperlinks(doc):
    hyperlinks = []
    rels = doc.part.rels
    for rel in rels:
        if rels[rel].reltype == RT.HYPERLINK:
            hyperlinks.append(rels[rel]._target)
    return hyperlinks     

def names_in_table(tables):
    names = []
    for table in tables:
        #Look if in the headers there is a synonym for author
        for i, cell in enumerate(table.rows[0].cells):
            if cell.text.lower() in author_synonyms:
                #Looking for names in the rows
                for row in table.rows[1:]:
                    names.append(row.cells[i].text)
    return names

def get_headers(doc):
    headers = []
    section = doc.sections[0]
    header = section.header
    for paragraph in header.paragraphs:
        if paragraph.text not in headers:
            headers.append(paragraph.text)
    return headers if headers != [''] else []

def get_footers(doc):
    footers = []
    section = doc.sections[0]
    footer = section.footer
    for paragraph in footer.paragraphs:
        if paragraph.text not in footers:
            footers.append(paragraph.text)
    return footers if footers != [''] else []

def path_to_filename(path):
    """
    Return the file name from a path
    """
    return os.path.splitext(os.path.basename(path))[0]